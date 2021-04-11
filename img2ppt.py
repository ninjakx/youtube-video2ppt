from pptx import Presentation  
from pptx.util import Inches 
import pptx 
from pathlib import Path
import argparse
import re
from glob import glob

def natural_sort_key(s, _nsre=re.compile('([0-9]+)')):
    return [
        int(text)
        if text.isdigit() else text.lower()
        for text in _nsre.split(s)]

# ap = argparse.ArgumentParser()
# ap.add_argument("-p", "--path", required=False, default='', help="Path to load img files")
# args = vars(ap.parse_args())


def create_ppt(img_path):
    prs = Presentation()  
    blank_slide_layout = prs.slide_layouts[6]  

    img_path = Path(img_path)
    # img_path = Path(args['path'])

    images = [img.name for img in img_path.iterdir() if img.suffix==".jpg"]

    sorted_images = sorted(images, key=natural_sort_key)

    for img in sorted_images:
        slide = prs.slides.add_slide(blank_slide_layout) 
        pic = slide.shapes.add_picture(f"{img_path/Path(img)}", Inches(0.5), Inches(0.75), width=Inches(9.2), height=Inches(6)) 

    [f.unlink() for f in img_path.glob("*") if f.is_file()] 
    prs.save(img_path/Path('slide_notes.pptx'))


